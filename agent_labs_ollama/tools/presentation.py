#!/usr/bin/env python3
"""
Presentation Tool
Generates PowerPoint presentations from structured slide data
"""

import json
import sys
import os
import base64
from typing import Dict, Any, List
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

def create_presentation(slides_data: List[Dict[str, Any]], title: str = "", output_filename: str = "presentation.pptx") -> Dict[str, Any]:
    """
    Create PowerPoint presentation from structured slide data

    Args:
        slides_data: List of slide dictionaries with 'title' and 'content' keys
        title: Presentation title
        output_filename: Output file name

    Returns:
        Dictionary with creation results
    """
    try:
        if not PPTX_AVAILABLE:
            return {
                "tool": "presentation",
                "success": False,
                "error": "python-pptx library not available. Install with: pip install python-pptx"
            }

        if not slides_data:
            return {
                "tool": "presentation",
                "success": False,
                "error": "No slide data provided"
            }

        # Create PowerPoint presentation
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title or "Generated Presentation"
        slide.shapes.placeholders[1].text = f"Generated on {datetime.now().strftime('%Y-%m-%d')}"

        slide_descriptions = []

        # Process each slide from structured data
        for i, slide_data in enumerate(slides_data, 1):
            slide_title = slide_data.get("title", f"Slide {i}")
            slide_content = slide_data.get("content", [])

            # Ensure content is a list
            if isinstance(slide_content, str):
                slide_content = [slide_content]
            elif not isinstance(slide_content, list):
                slide_content = [str(slide_content)]

            # Add slide to presentation
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            slide.shapes.title.text = slide_title

            content_placeholder = slide.shapes.placeholders[1]
            text_frame = content_placeholder.text_frame
            text_frame.clear()

            for j, bullet in enumerate(slide_content[:6]):  # Limit to 6 bullets per slide
                if j == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = str(bullet)
                p.level = 0

            slide_descriptions.append({
                "slide_number": len(prs.slides),
                "title": slide_title,
                "bullets": slide_content,  # Frontend expects 'bullets' not 'content'
                "slide_type": "title" if i == 1 else "content"  # First content slide is title slide
            })

        # Save presentation to outputs directory
        outputs_dir = os.path.join(os.path.dirname(__file__), "..", "outputs")
        os.makedirs(outputs_dir, exist_ok=True)
        output_path = os.path.join(outputs_dir, output_filename)
        prs.save(output_path)

        # Convert to base64 for JSON response
        with open(output_path, "rb") as f:
            ppt_data = f.read()
            ppt_base64 = base64.b64encode(ppt_data).decode('utf-8')

        # Calculate file size in MB
        file_size_mb = len(ppt_data) / (1024 * 1024)

        return {
            "tool": "presentation",
            "success": True,
            "title": title,
            "slides_created": len(prs.slides),  # Frontend expects 'slides_created'
            "total_slides": len(prs.slides),     # Frontend expects 'total_slides'
            "file_size_mb": round(file_size_mb, 2),  # Frontend expects 'file_size_mb'
            "output_path": output_path,
            "presentation_data": {
                "base64": ppt_base64,
                "filename": output_filename,
                "mime_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "slides": slide_descriptions
            },
            "timestamp": datetime.now().isoformat()
        }

    except Exception as e:
        return {
            "tool": "presentation",
            "success": False,
            "error": str(e)
        }

def main():
    """CLI interface for the presentation tool"""
    if len(sys.argv) < 2:
        print(json.dumps({"error": "Usage: presentation.py <json_args>"}))
        sys.exit(1)

    try:
        args = json.loads(sys.argv[1])
        slides_data = args.get("slides_data", [])
        title = args.get("title", "Generated Presentation")
        output_filename = args.get("output_filename", "presentation.pptx")

        if not slides_data:
            print(json.dumps({"error": "slides_data is required"}))
            sys.exit(1)

        result = create_presentation(slides_data, title, output_filename)
        print(json.dumps(result, indent=2))

    except json.JSONDecodeError as e:
        print(json.dumps({"error": f"Invalid JSON arguments: {e}"}))
        sys.exit(1)
    except Exception as e:
        print(json.dumps({"error": str(e)}))
        sys.exit(1)

if __name__ == "__main__":
    main()