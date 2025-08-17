import sys
import asyncio
import PyPDF2
import aiohttp
import re
import os
import subprocess
import platform
from pathlib import Path
from pptx import Presentation
from typing import List, Dict, Any

from PyQt6.QtWidgets import (
    QApplication, QWidget, QVBoxLayout, QHBoxLayout, QPushButton, QLabel,
    QFileDialog, QTextEdit, QProgressBar, QLineEdit, QComboBox, QFrame,
    QGroupBox, QSizePolicy, QSpacerItem, QScrollArea
)
from PyQt6.QtCore import Qt, QThread, pyqtSignal, QPropertyAnimation, QEasingCurve, QRect
from PyQt6.QtGui import QFont, QPalette, QColor, QPixmap, QPainter, QBrush, QLinearGradient

def log_progress(msg, progress_cb=None):
    if progress_cb:
        progress_cb(msg)
    else:
        print(msg)

def open_file_location(file_path):
    """Open the folder containing the specified file in the system file manager"""
    try:
        file_path = Path(file_path).resolve()
        folder_path = file_path.parent
        
        system = platform.system()
        if system == "Darwin":  # macOS
            subprocess.run(["open", "-R", str(file_path)], check=True)
        elif system == "Windows":
            subprocess.run(["explorer", "/select,", str(file_path)], check=True)
        else:  # Linux and other Unix-like systems
            # Try common file managers
            file_managers = ["nautilus", "dolphin", "thunar", "caja", "pcmanfm"]
            opened = False
            for fm in file_managers:
                try:
                    subprocess.run([fm, str(folder_path)], check=True)
                    opened = True
                    break
                except (subprocess.CalledProcessError, FileNotFoundError):
                    continue
            
            if not opened:
                # Fallback: try to open with xdg-open
                subprocess.run(["xdg-open", str(folder_path)], check=True)
        
        return True
    except Exception as e:
        print(f"Warning: Could not open file location: {e}")
        return False

class OllamaAPI:
    """Handle communication with Ollama API"""
    
    def __init__(self, base_url: str = "http://127.0.0.1:11434"):
        self.base_url = base_url
        self.session = None
    
    async def __aenter__(self):
        self.session = aiohttp.ClientSession()
        return self
    
    async def __aexit__(self, exc_type, exc_val, exc_tb):
        if self.session:
            await self.session.close()
    
    async def get_models(self) -> List[Dict[str, Any]]:
        """Get available Ollama models"""
        try:
            timeout = aiohttp.ClientTimeout(total=30)
            async with self.session.get(f"{self.base_url}/api/tags", timeout=timeout) as response:
                if response.status == 200:
                    data = await response.json()
                    models = []
                    for model in data.get('models', []):
                        models.append({
                            'name': model['name'],
                            'display_name': model['name'].replace(':latest', ''),
                            'size': model.get('size', 0),
                            'family': model.get('details', {}).get('family', 'unknown')
                        })
                    return models
                else:
                    print(f"Failed to fetch models: HTTP {response.status}")
                    return []
        except Exception as e:
            print(f"Error getting models: {e}")
            return []

class ModelLoader(QThread):
    """Worker thread for loading models"""
    models_loaded = pyqtSignal(list)
    error_occurred = pyqtSignal(str)
    
    def run(self):
        """Load models in background"""
        try:
            loop = asyncio.new_event_loop()
            asyncio.set_event_loop(loop)
            
            async def load_models():
                async with OllamaAPI() as api:
                    models = await api.get_models()
                    return models
            
            models = loop.run_until_complete(load_models())
            loop.close()
            
            self.models_loaded.emit(models)
            
        except Exception as e:
            self.error_occurred.emit(str(e))

async def process_pdf(pdf_path, model_name='gemma3:latest', slide_layout=1, progress_cb=None):
    try:
        with open(pdf_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            text = ''.join(page.extract_text() or '' for page in reader.pages)
    except Exception as e:
        log_progress(f'❌ Failed to open PDF: {e}', progress_cb)
        return

    # Print the extracted text for debugging (first 500 chars)
    debug_text = text[:500] + ("..." if len(text) > 500 else "")
    log_progress(f'--- Extracted PDF text (first 500 chars) ---\n{debug_text}\n--- End of extract ---', progress_cb)

    if not text.strip():
        log_progress('⚠️ Warning: Extracted PDF text is empty.', progress_cb)
        return

    # Chunking utility
    def chunk_text(text, max_length=3000):
        """Splits text into chunks of max_length characters, trying to preserve paragraphs."""
        # Split text into chunks of max_length characters, on paragraph boundaries if possible
        paragraphs = text.split('\n')
        chunks = []
        current = ''
        for para in paragraphs:
            if len(current) + len(para) + 1 > max_length:
                if current:
                    chunks.append(current)
                current = para
            else:
                current += ('\n' if current else '') + para
        if current:
            chunks.append(current)
        return chunks

    # If text is too long, chunk it and process each chunk separately
    max_chunk_length = 3000  # adjust as needed for your LLM
    text_chunks = chunk_text(text, max_chunk_length)
    all_slides = []
    for chunk_idx, chunk in enumerate(text_chunks):
        chunk_prompt = f'''
You are a presentation assistant. Your ONLY task is to convert the following document into a slide outline.

IMPORTANT: DO NOT provide any commentary, review, explanation, or text before or after the slides. If you do, the user's program will break.

INSTRUCTIONS:
- Output ONLY slides in the following format. Do NOT add any extra text, commentary, or explanation.
- Always generate at least 3 slides.
- Each slide must have a title and 3 to 5 bullet points.
- Use this exact format (do not change it):

Slide 1:
Title: <title of slide 1>
- <bullet point 1>
- <bullet point 2>
- <bullet point 3>
Slide 2:
Title: <title of slide 2>
- <bullet point 1>
- <bullet point 2>
- <bullet point 3>
Slide 3:
Title: <title of slide 3>
- <bullet point 1>
- <bullet point 2>
- <bullet point 3>

Continue for as many slides as needed, but do not add any text outside this format. Do NOT include any summary, review, or extra lines.

Document:
{chunk}
'''
        log_progress(f'⏳ Sending request to LLM for chunk {chunk_idx+1}/{len(text_chunks)}...', progress_cb)
        try:
            async with aiohttp.ClientSession() as session:
                async with session.post('http://localhost:11434/api/generate', json={
                    'model': model_name,
                    'prompt': chunk_prompt,
                    'stream': False
                }) as resp:
                    result = await resp.json()
        except Exception as e:
            log_progress(f'❌ LLM API error: {e}', progress_cb)
            return

        if 'response' not in result:
            log_progress(f'⚠️ Unexpected API response: {result}', progress_cb)
            return

        outline = result['response']
        log_progress(f'=== LLM Output Start (chunk {chunk_idx+1}) ===\n' + outline + f'\n=== LLM Output End (chunk {chunk_idx+1}) ===', progress_cb)

        slides_raw = re.split(r'Slide\s*\d+:', outline, flags=re.IGNORECASE)[1:]
        if not slides_raw or all(not s.strip() for s in slides_raw):
            log_progress((f"❌ LLM output for chunk {chunk_idx+1} did not contain any slides in the expected format. "
                   "Please check your LLM prompt and model output.\n"
                   "First 500 chars of LLM output:\n" + outline[:500] + ("..." if len(outline) > 500 else "")), progress_cb)
            continue  # skip this chunk
        all_slides.extend(slides_raw)

    if not all_slides:
        log_progress(("❌ No valid slides were generated from any chunk. "
               "Please check your LLM prompt, model, and PDF content."), progress_cb)
        return

    prs = Presentation()
    for idx, slide_text in enumerate(all_slides, 1):
        lines = [line.strip() for line in slide_text.strip().splitlines() if line.strip()]
        title = None
        bullets = []
        for i, line in enumerate(lines):
            if line.lower().startswith('title:'):
                title = line[len('title:'):].strip()
                bullets = [l[1:].strip() for l in lines[i+1:] if l.startswith('-')]
                break
        if not title:
            if lines:
                title = lines[0]
                bullets = [l[1:].strip() for l in lines[1:] if l.startswith('-')]
            else:
                title = f"Slide {idx}"
                bullets = []
        slide = prs.slides.add_slide(prs.slide_layouts[slide_layout])
        slide.shapes.title.text = title
        slide.placeholders[1].text = '\n'.join(bullets)
        log_progress(f'✅ Added Slide {idx}: "{title}" with {len(bullets)} bullet points.', progress_cb)

    pdf_name = os.path.splitext(os.path.basename(pdf_path))[0]
    output_file = f'{pdf_name}_presentation.pptx'
    
    # Save the presentation
    prs.save(output_file)
    
    # Get the absolute path for better user feedback
    abs_output_path = os.path.abspath(output_file)
    
    log_progress(f'🎉 Presentation saved as {abs_output_path}', progress_cb)
    
    # Open the folder containing the generated file
    if open_file_location(abs_output_path):
        log_progress(f'📂 Opened folder: {os.path.dirname(abs_output_path)}', progress_cb)
    else:
        log_progress(f'📁 File saved to: {os.path.dirname(abs_output_path)}', progress_cb)
    
    return abs_output_path
    # ...existing code...
class Worker(QThread):
    progress = pyqtSignal(str)
    finished = pyqtSignal(str)

    def __init__(self, pdf_path, model_name, slide_layout):
        super().__init__()
        self.pdf_path = pdf_path
        self.model_name = model_name
        self.slide_layout = slide_layout
        self.task = None

    def run(self):
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        self.task = loop.create_task(self._run())
        try:
            loop.run_until_complete(self.task)
        except asyncio.CancelledError:
            pass
        finally:
            loop.close()

    def stop(self):
        if self.task and not self.task.done():
            self.task.cancel()
        self.quit()

    async def _run(self):
        try:
            await process_pdf(self.pdf_path, self.model_name, self.slide_layout, self.progress.emit)
            self.finished.emit("Done")
        except asyncio.CancelledError:
            self.finished.emit("Cancelled")

class MainWindow(QWidget):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("PDF to PPTX Converter")
        self.setGeometry(300, 300, 800, 800)
        self.setMinimumSize(700, 700)
        
        # Initialize model state
        self.available_models = []
        self.current_model = None
        self.worker = None
        
        # Apply modern styling
        self.setup_styles()
        self.setup_ui()
        
        # Load models on startup
        self.load_models()

    def setup_styles(self):
        """Apply modern styling to the application"""
        self.setStyleSheet("""
            QWidget {
                background-color: #f8f9fa;
                color: #2c3e50;
                font-family: 'Helvetica Neue', Arial;
            }
            
            QGroupBox {
                font-size: 14px;
                font-weight: 600;
                color: #34495e;
                border: 2px solid #e9ecef;
                border-radius: 12px;
                margin-top: 10px;
                padding-top: 15px;
                background-color: white;
            }
            
            QGroupBox::title {
                subcontrol-origin: margin;
                left: 15px;
                padding: 0 8px 0 8px;
                background-color: white;
            }
            
            QPushButton {
                background: qlineargradient(x1:0, y1:0, x2:0, y2:1,
                    stop:0 #3498db, stop:1 #2980b9);
                border: none;
                border-radius: 8px;
                color: white;
                font-size: 13px;
                font-weight: 600;
                padding: 12px 24px;
                min-height: 20px;
            }
            
            QPushButton:hover {
                background: qlineargradient(x1:0, y1:0, x2:0, y2:1,
                    stop:0 #5dade2, stop:1 #3498db);
            }
            
            QPushButton:pressed {
                background: qlineargradient(x1:0, y1:0, x2:0, y2:1,
                    stop:0 #2980b9, stop:1 #1f618d);
            }
            
            QPushButton:disabled {
                background: #bdc3c7;
                color: #7f8c8d;
            }
            
            QPushButton#primaryButton {
                background: qlineargradient(x1:0, y1:0, x2:0, y2:1,
                    stop:0 #27ae60, stop:1 #229954);
                font-size: 14px;
                padding: 14px 28px;
            }
            
            QPushButton#primaryButton:hover {
                background: qlineargradient(x1:0, y1:0, x2:0, y2:1,
                    stop:0 #58d68d, stop:1 #27ae60);
            }
            
            QPushButton#dangerButton {
                background: qlineargradient(x1:0, y1:0, x2:0, y2:1,
                    stop:0 #e74c3c, stop:1 #c0392b);
            }
            
            QPushButton#dangerButton:hover {
                background: qlineargradient(x1:0, y1:0, x2:0, y2:1,
                    stop:0 #ec7063, stop:1 #e74c3c);
            }
            
            QPushButton#refreshButton {
                background: qlineargradient(x1:0, y1:0, x2:0, y2:1,
                    stop:0 #f39c12, stop:1 #e67e22);
                padding: 8px 16px;
                font-size: 12px;
            }
            
            QPushButton#refreshButton:hover {
                background: qlineargradient(x1:0, y1:0, x2:0, y2:1,
                    stop:0 #f7dc6f, stop:1 #f39c12);
            }
            
            QComboBox {
                border: 2px solid #e9ecef;
                border-radius: 8px;
                padding: 10px 15px;
                background-color: white;
                font-size: 13px;
                min-height: 20px;
            }
            
            QComboBox:focus {
                border-color: #3498db;
            }
            
            QComboBox::drop-down {
                border: none;
                width: 30px;
            }
            
            QComboBox::down-arrow {
                image: none;
                border: 4px solid transparent;
                border-top: 6px solid #7f8c8d;
                margin-right: 10px;
            }
            
            QTextEdit {
                border: 2px solid #e9ecef;
                border-radius: 8px;
                background-color: white;
                padding: 15px 15px 25px 15px;
                font-family: Monaco, 'Courier New';
                font-size: 12px;
                line-height: 1.5;
            }
            
            QTextEdit:focus {
                border-color: #3498db;
            }
            
            QProgressBar {
                border: none;
                border-radius: 10px;
                background-color: #ecf0f1;
                height: 20px;
                text-align: center;
            }
            
            QProgressBar::chunk {
                background: qlineargradient(x1:0, y1:0, x2:1, y2:0,
                    stop:0 #3498db, stop:1 #2980b9);
                border-radius: 10px;
            }
            
            QLabel {
                color: #2c3e50;
                font-size: 13px;
            }
            
            QLabel#headerLabel {
                font-size: 18px;
                font-weight: 700;
                color: #2c3e50;
                margin: 10px 0;
            }
            
            QLabel#statusLabel {
                font-size: 14px;
                color: #7f8c8d;
                margin: 5px 0;
            }
        """)

    def setup_ui(self):
        """Setup the modern UI layout"""
        main_layout = QVBoxLayout()
        main_layout.setSpacing(15)
        main_layout.setContentsMargins(20, 20, 20, 20)
        self.setLayout(main_layout)
        
        # Header section
        header_label = QLabel("PDF to PowerPoint Converter")
        header_label.setObjectName("headerLabel")
        header_label.setAlignment(Qt.AlignmentFlag.AlignCenter)
        main_layout.addWidget(header_label)
        
        # Status label
        self.status_label = QLabel("Ready to convert your PDF files to presentations")
        self.status_label.setObjectName("statusLabel")
        self.status_label.setAlignment(Qt.AlignmentFlag.AlignCenter)
        main_layout.addWidget(self.status_label)
        
        # Configuration section
        config_group = QGroupBox("Configuration")
        config_layout = QVBoxLayout()
        config_layout.setSpacing(15)
        config_group.setLayout(config_layout)
        
        # Model selection section
        model_section = QFrame()
        model_layout = QVBoxLayout()
        model_layout.setContentsMargins(0, 0, 0, 0)
        model_section.setLayout(model_layout)
        
        model_label = QLabel("Ollama Model:")
        model_label.setStyleSheet("font-weight: 600; margin-bottom: 5px;")
        model_layout.addWidget(model_label)
        
        model_controls_layout = QHBoxLayout()
        model_controls_layout.setSpacing(10)
        
        self.model_combo = QComboBox()
        self.model_combo.setMinimumWidth(300)
        self.model_combo.addItem("Loading models...")
        model_controls_layout.addWidget(self.model_combo)
        
        self.refresh_models_btn = QPushButton("🔄 Refresh")
        self.refresh_models_btn.setObjectName("refreshButton")
        self.refresh_models_btn.clicked.connect(self.load_models)
        model_controls_layout.addWidget(self.refresh_models_btn)
        
        model_controls_layout.addStretch()
        model_layout.addLayout(model_controls_layout)
        config_layout.addWidget(model_section)
        
        # Layout selection section
        layout_section = QFrame()
        layout_layout = QVBoxLayout()
        layout_layout.setContentsMargins(0, 0, 0, 0)
        layout_section.setLayout(layout_layout)
        
        layout_label = QLabel("Slide Layout:")
        layout_label.setStyleSheet("font-weight: 600; margin-bottom: 5px;")
        layout_layout.addWidget(layout_label)
        
        self.layout_combo = QComboBox()
        self.layout_combo.addItems([
            "0: Title Slide",
            "1: Title and Content",
            "2: Section Header", 
            "3: Two Content",
            "4: Comparison",
            "5: Title Only",
            "6: Blank",
            "7: Content with Caption",
            "8: Picture with Caption"
        ])
        self.layout_combo.setCurrentIndex(1)
        layout_layout.addWidget(self.layout_combo)
        config_layout.addWidget(layout_section)
        
        main_layout.addWidget(config_group)
        
        # Action buttons section
        button_layout = QHBoxLayout()
        button_layout.setSpacing(15)
        
        self.select_btn = QPushButton("📁 Select PDF File")
        self.select_btn.setObjectName("primaryButton")
        self.select_btn.clicked.connect(self.select_pdf)
        button_layout.addWidget(self.select_btn)
        
        self.cancel_btn = QPushButton("❌ Cancel Processing")
        self.cancel_btn.setObjectName("dangerButton")
        self.cancel_btn.clicked.connect(self.cancel_processing)
        self.cancel_btn.setEnabled(False)
        button_layout.addWidget(self.cancel_btn)
        
        button_layout.addStretch()
        main_layout.addLayout(button_layout)
        
        # Progress section
        progress_group = QGroupBox("Progress")
        progress_layout = QVBoxLayout()
        progress_group.setLayout(progress_layout)
        
        self.progress_bar = QProgressBar()
        self.progress_bar.setRange(0, 0)
        self.progress_bar.setVisible(False)
        progress_layout.addWidget(self.progress_bar)
        
        # Output section with proper text display
        self.text_edit = QTextEdit()
        self.text_edit.setReadOnly(True)
        self.text_edit.setPlaceholderText("Processing logs will appear here...")
        self.text_edit.setMinimumHeight(200)
        
        # Configure scrolling for better visibility
        self.text_edit.setHorizontalScrollBarPolicy(Qt.ScrollBarPolicy.ScrollBarAsNeeded)
        self.text_edit.setVerticalScrollBarPolicy(Qt.ScrollBarPolicy.ScrollBarAsNeeded)
        
        progress_layout.addWidget(self.text_edit)
        
        main_layout.addWidget(progress_group)

    def load_models(self):
        """Load available Ollama models"""
        self.refresh_models_btn.setText("⏳ Loading...")
        self.refresh_models_btn.setEnabled(False)
        self.model_combo.clear()
        self.model_combo.addItem("Loading models...")
        self.status_label.setText("Loading available Ollama models...")
        
        self.model_loader = ModelLoader()
        self.model_loader.models_loaded.connect(self.on_models_loaded)
        self.model_loader.error_occurred.connect(self.on_models_error)
        self.model_loader.start()
    
    def on_models_loaded(self, models):
        """Handle loaded models"""
        self.available_models = models
        self.model_combo.clear()
        
        if not models:
            self.model_combo.addItem("❌ No models available")
            self.refresh_models_btn.setText("🔄 Refresh")
            self.refresh_models_btn.setEnabled(True)
            self.status_label.setText("No Ollama models found. Please install models first.")
            return
        
        # Filter out embedding models
        chat_models = [m for m in models if 'embed' not in m['family'] and 'embed' not in m['name']]
        
        if not chat_models:
            self.model_combo.addItem("❌ No chat models available")
            self.refresh_models_btn.setText("🔄 Refresh")
            self.refresh_models_btn.setEnabled(True)
            self.status_label.setText("No suitable chat models found.")
            return
        
        # Populate combo box
        for model in chat_models:
            display_name = model['display_name']
            size_gb = model['size'] / (1024 * 1024 * 1024)
            display_text = f"{display_name} ({size_gb:.1f}GB)"
            self.model_combo.addItem(display_text, model)
        
        # Set default model (prefer gemma3, then first available)
        default_index = 0
        for i, model in enumerate(chat_models):
            if 'gemma' in model['name'].lower():
                default_index = i
                break
        
        self.model_combo.setCurrentIndex(default_index)
        self.current_model = chat_models[default_index] if chat_models else None
        
        # Re-enable refresh button
        self.refresh_models_btn.setText("🔄 Refresh")
        self.refresh_models_btn.setEnabled(True)
        
        if self.current_model:
            self.status_label.setText(f"Ready to convert PDFs using {self.current_model['display_name']}")
    
    def on_models_error(self, error):
        """Handle model loading error"""
        self.model_combo.clear()
        self.model_combo.addItem("❌ Connection failed")
        self.refresh_models_btn.setText("🔄 Refresh")
        self.refresh_models_btn.setEnabled(True)
        self.status_label.setText("Failed to connect to Ollama. Please ensure Ollama is running.")

    def select_pdf(self):
        """Select and process PDF file"""
        file_path, _ = QFileDialog.getOpenFileName(
            self, 
            "Select PDF File", 
            "", 
            "PDF Files (*.pdf);;All Files (*)"
        )
        if file_path:
            # Validate file exists and is readable
            if not os.path.exists(file_path):
                self.status_label.setText("❌ Selected file does not exist.")
                self.text_edit.append(f"❌ Error: File not found: {file_path}")
                return
            
            if not os.path.isfile(file_path):
                self.status_label.setText("❌ Selected path is not a file.")
                self.text_edit.append(f"❌ Error: Path is not a file: {file_path}")
                return
            
            if not os.access(file_path, os.R_OK):
                self.status_label.setText("❌ Cannot read the selected file.")
                self.text_edit.append(f"❌ Error: No read permission for file: {file_path}")
                return
            
            # Check file extension
            if not file_path.lower().endswith('.pdf'):
                self.status_label.setText("⚠️ Selected file may not be a PDF.")
                self.text_edit.append(f"⚠️ Warning: File does not have .pdf extension: {file_path}")
            
            # Check if a model is selected
            current_model_data = self.model_combo.currentData()
            if not current_model_data:
                self.status_label.setText("❌ No model selected. Please select a model first.")
                return
            
            # Clear previous output and setup UI for processing
            self.text_edit.clear()
            self.progress_bar.setVisible(True)
            self.status_label.setText(f"🔄 Processing: {os.path.basename(file_path)}")
            
            # Add file info to output
            file_size = os.path.getsize(file_path) / (1024 * 1024)  # MB
            self.text_edit.append(f"📄 Selected file: {os.path.basename(file_path)}")
            self.text_edit.append(f"📐 File size: {file_size:.1f} MB")
            self.text_edit.append(f"📍 Full path: {file_path}")
            self.text_edit.append("─" * 50)
            
            # Start processing
            model_name = current_model_data['name']
            slide_layout = self.layout_combo.currentIndex()
            self.worker = Worker(file_path, model_name, slide_layout)
            self.worker.progress.connect(self.append_text)
            self.worker.finished.connect(self.done)
            self.worker.start()
            
            # Update UI state
            self.select_btn.setEnabled(False)
            self.cancel_btn.setEnabled(True)
            self.select_btn.setText("🔄 Processing...")

    def cancel_processing(self):
        """Cancel the current processing task"""
        if self.worker and self.worker.isRunning():
            self.worker.stop()
            self.cancel_btn.setEnabled(False)
            self.select_btn.setEnabled(True)
            self.select_btn.setText("📁 Select PDF File")
            self.progress_bar.setVisible(False)
            self.status_label.setText("❌ Processing cancelled")
            self.text_edit.append("\n❌ Processing was cancelled by user")

    def append_text(self, msg):
        """Append text to the output area and auto-scroll"""
        # Format long lines for better readability
        formatted_msg = self.format_progress_message(msg)
        self.text_edit.append(formatted_msg)
        # Auto-scroll to bottom with a small delay to ensure proper rendering
        self.text_edit.ensureCursorVisible()
        scrollbar = self.text_edit.verticalScrollBar()
        scrollbar.setValue(scrollbar.maximum())
        # Add an extra empty line for better visibility of the last message
        cursor = self.text_edit.textCursor()
        cursor.movePosition(cursor.MoveOperation.End)
        self.text_edit.setTextCursor(cursor)
    
    def format_progress_message(self, msg):
        """Format progress messages for better display"""
        # If message is very long (like file paths), add line breaks for readability
        if len(msg) > 100:
            # Look for file paths and break them nicely
            if "saved as" in msg.lower() or "full path:" in msg.lower():
                # Find the path part and format it
                parts = msg.split(": ", 1)
                if len(parts) == 2:
                    prefix, path = parts
                    # Break long paths into multiple lines if needed
                    if len(path) > 80:
                        return f"{prefix}:\n    {path}"
                    else:
                        return msg
        return msg

    def done(self, msg):
        """Handle processing completion"""
        self.progress_bar.setVisible(False)
        self.select_btn.setEnabled(True)
        self.cancel_btn.setEnabled(False)
        self.select_btn.setText("📁 Select PDF File")
        
        if msg == "Done":
            self.status_label.setText("✅ Conversion completed successfully!")
            self.text_edit.append("\n🎉 Conversion completed! Your PowerPoint file has been saved and the folder has been opened.")
            self.text_edit.append("\n💡 Tip: You can now open the PowerPoint file to view your presentation!")
        elif msg == "Cancelled":
            self.status_label.setText("❌ Processing was cancelled")
            self.text_edit.append("\n❌ Processing was cancelled")
        else:
            self.status_label.setText("❌ Processing failed")
            self.text_edit.append(f"\n❌ Processing failed: {msg}")
            self.text_edit.append("\n💡 Tip: Check if Ollama is running and the selected model is available.")
        
        # Auto-scroll to bottom and ensure cursor visibility
        self.text_edit.ensureCursorVisible()
        scrollbar = self.text_edit.verticalScrollBar()
        scrollbar.setValue(scrollbar.maximum())
        cursor = self.text_edit.textCursor()
        cursor.movePosition(cursor.MoveOperation.End)
        self.text_edit.setTextCursor(cursor)

def main():
    # CLI mode
    if '--pdf_file' in sys.argv:
        idx = sys.argv.index('--pdf_file')
        if idx + 1 < len(sys.argv):
            pdf_path = sys.argv[idx + 1]
            model_name = 'gemma3:latest'
            if '--model' in sys.argv:
                model_idx = sys.argv.index('--model')
                if model_idx + 1 < len(sys.argv):
                    model_name = sys.argv[model_idx + 1]
            slide_layout = 1
            if '--slide_layout' in sys.argv:
                layout_idx = sys.argv.index('--slide_layout')
                if layout_idx + 1 < len(sys.argv):
                    try:
                        slide_layout = int(sys.argv[layout_idx + 1])
                    except ValueError:
                        print("Invalid slide layout, using default.")
            asyncio.run(process_pdf(pdf_path, model_name=model_name, slide_layout=slide_layout))
        else:
            print("Usage: python gen_ppt.py --pdf_file <path_to_pdf> [--model <model_name>] [--slide_layout <layout_index>]")
        return

    # GUI mode
    app = QApplication(sys.argv)
    window = MainWindow()
    window.show()
    sys.exit(app.exec())

if __name__ == '__main__':
    main()
